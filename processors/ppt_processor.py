from pptx import Presentation
from typing import Dict, Any, List
from .base_processor import BaseProcessor

class PPTProcessor(BaseProcessor):
    """Process PowerPoint files"""
    
    def process(self, filepath: str) -> Dict[str, Any]:
        """Extract text and structure from PowerPoint presentations"""
        result = {
            'text': '',
            'slides': [],
            'metadata': {}
        }
        
        try:
            prs = Presentation(filepath)
            
            # Extract metadata
            core_props = prs.core_properties
            result['metadata'] = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'slide_count': len(prs.slides)
            }
            
            text_parts = [f"# Presentation\n\n**Slides:** {len(prs.slides)}\n\n"]
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': slide_idx,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                slide_text = []
                slide_text.append(f"## Slide {slide_idx}\n\n")
                
                # Extract text from shapes (including grouped shapes)
                for shape in slide.shapes:
                    self._extract_text_from_shape(shape, slide_data, slide_text)
                
                # Extract speaker notes (improved method)
                slide_data['notes'] = self._extract_speaker_notes(slide)
                if slide_data['notes']:
                    slide_text.append(f"**Speaker Notes:** {slide_data['notes']}\n\n")
                
                result['slides'].append(slide_data)
                text_parts.extend(slide_text)
            
            result['text'] = '\n'.join(text_parts)
            
        except Exception as e:
            result['error'] = str(e)
            result['text'] = f"Error processing PowerPoint: {str(e)}"
        
        return result
    
    def to_markdown(self, content: Dict[str, Any]) -> str:
        """Convert PowerPoint content to markdown"""
        if isinstance(content, str):
            return content
        
        if content.get('text'):
            return content['text']
        
        # Fallback conversion
        markdown = "# Presentation\n\n"
        
        # Add metadata
        if content.get('metadata'):
            metadata = content['metadata']
            markdown += "## Information\n\n"
            for key, value in metadata.items():
                if value:
                    markdown += f"**{key.title()}:** {value}\n\n"
        
        # Add slides
        if content.get('slides'):
            markdown += "## Slides\n\n"
            for slide in content['slides']:
                markdown += f"### Slide {slide['slide_number']}"
                
                if slide.get('title'):
                    markdown += f": {slide['title']}\n\n"
                else:
                    markdown += "\n\n"
                
                if slide.get('content'):
                    for content_item in slide['content']:
                        markdown += f"{content_item}\n\n"
                
                if slide.get('notes'):
                    markdown += f"**Speaker Notes:** {slide['notes']}\n\n"
                
                markdown += "---\n\n"
        
        return markdown
    
    def _extract_speaker_notes(self, slide) -> str:
        """Enhanced method to extract speaker notes from a slide"""
        notes_text = ""
        
        try:
            # Check if slide has notes
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_slide = slide.notes_slide
                
                # Method 1: Try to get notes from notes_text_frame (most reliable)
                if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                
                # Method 2: Fallback - iterate through shapes to find notes
                if not notes_text:
                    for shape in notes_slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            # Skip placeholder text
                            shape_text = shape.text.strip()
                            if shape_text and shape_text.lower() not in ['click to add notes', 'notes']:
                                notes_text = shape_text
                                break
                
                # Method 3: Alternative attribute check for older python-pptx versions
                if not notes_text and hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                    try:
                        notes_slide = slide.notes_slide
                        if notes_slide.shapes:
                            for shape in notes_slide.shapes:
                                if hasattr(shape, 'text_frame') and shape.text_frame:
                                    text = shape.text_frame.text.strip()
                                    if text and text.lower() not in ['click to add notes', 'notes']:
                                        notes_text = text
                                        break
                    except:
                        pass  # Ignore errors in fallback method
                        
        except Exception as e:
            # Log error but don't fail the entire processing
            print(f"Warning: Could not extract speaker notes: {str(e)}")
        
        return notes_text
    
    def _extract_text_from_shape(self, shape, slide_data, slide_text):
        """Extract text from a shape, including grouped shapes"""
        try:
            # Check if this is a group shape
            if hasattr(shape, 'shape_type') and shape.shape_type == 6:  # Group shape type
                # Recursively extract text from each shape in the group
                if hasattr(shape, 'shapes'):
                    for grouped_shape in shape.shapes:
                        self._extract_text_from_shape(grouped_shape, slide_data, slide_text)
            elif hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                
                # Try to identify title
                if not slide_data['title'] and (
                    shape.shape_type == 14 or  # Title placeholder
                    len(text.split('\n')[0]) < 100  # Short first line likely a title
                ):
                    slide_data['title'] = text.split('\n')[0]
                    slide_text.append(f"### {slide_data['title']}\n\n")
                    
                    # Add remaining content if any
                    remaining_text = '\n'.join(text.split('\n')[1:]).strip()
                    if remaining_text:
                        slide_data['content'].append(remaining_text)
                        slide_text.append(f"{remaining_text}\n\n")
                else:
                    slide_data['content'].append(text)
                    slide_text.append(f"{text}\n\n")
        except Exception as e:
            # Log error but don't fail the entire processing
            print(f"Warning: Could not extract text from shape: {str(e)}")
